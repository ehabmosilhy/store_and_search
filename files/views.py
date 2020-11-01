from django.shortcuts import render, redirect
from .forms import FileForm
from .models import File
from django.http import HttpResponse
from django.conf import settings
from pptx import Presentation
import fitz
import os
from django.db.models import Q



def home(request):
    files = File.objects.all()
    return render(request, 'homepage.html', {
        'files': files
    })


def file_list(request):
    files = File.objects.all()
    return render(request, 'file_list.html', {
        'files': files
    })


def upload_file(request):
    if request.method == "POST":
        form = FileForm(request.POST, request.FILES)
        if form.is_valid():
            form.save()
            return redirect('file_list')
    else:
        form = FileForm()

    return render(request, 'upload_file.html',
                  {'form': form}
                  )


def make_search(keyword):
    prefix = "files/store"
    dir = os.path.join(settings.BASE_DIR, prefix)
    positive_files = []
    print(f'dir {dir}')
    files = os.listdir(dir)
    for file in files:
        if keyword in file:
            positive_files.append(f"{prefix}/{file}")
        else:
            myfile = os.path.join(dir, file)
            if '.pdf' in myfile:
                with open(myfile, "rb") as filehandle:
                    doc = fitz.open(myfile)
                    page1 = doc.loadPage(0)
                    page1text = page1.getText("text")

                    t = page1text
                    if keyword.lower() in t.lower():
                        positive_files.append(f"{prefix}/{file}")
            elif '.pptx' in myfile:
                prs = Presentation(myfile)
                found = 0
                for slide in prs.slides:
                    if found == 1:
                        break
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            print(shape.text.lower())
                            if keyword.lower() in shape.text.lower():
                                positive_files.append(f"{prefix}/{file}")
                                found = 1
                                break
    return positive_files


def search(request):
    search_text = request.POST['search_text']
    result = make_search(search_text)

    files = File.objects.filter(Q(content__in=result) | Q(title__contains=search_text) )

    return render(request, 'search.html', {'files': files})
